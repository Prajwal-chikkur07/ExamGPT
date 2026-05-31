"""Extract text from PDF / DOCX / PPTX / TXT and split into page-aware chunks."""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from langchain_text_splitters import RecursiveCharacterTextSplitter


@dataclass
class PageText:
    page: int
    text: str


@dataclass
class Chunk:
    page: int
    text: str


def _extract_pdf(path: Path) -> list[PageText]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages: list[PageText] = []
    for i, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        if text.strip():
            pages.append(PageText(page=i, text=text))

    # Image-only / scanned PDFs have no text layer — fall back to OCR.
    if not pages:
        from app.core.ocr import ocr_pdf
        import re

        raw = ocr_pdf(path)
        if raw.strip():
            parts = re.split(r"--- page (\d+) ---", raw)
            for i in range(1, len(parts), 2):
                try:
                    page_num = int(parts[i])
                except ValueError:
                    continue
                page_text = parts[i + 1].strip() if i + 1 < len(parts) else ""
                if page_text:
                    pages.append(PageText(page=page_num, text=page_text))
    return pages


def _extract_image(path: Path) -> list[PageText]:
    from app.core.ocr import ocr_image

    text = ocr_image(path)
    if not text.strip():
        return []
    return [PageText(page=1, text=text)]


def _extract_docx(path: Path) -> list[PageText]:
    from docx import Document

    doc = Document(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    # docx has no real pages — chunk by ~40 paragraphs as pseudo-pages.
    pages: list[PageText] = []
    page_size = 40
    for i in range(0, len(paragraphs), page_size):
        block = "\n".join(paragraphs[i : i + page_size])
        if block.strip():
            pages.append(PageText(page=(i // page_size) + 1, text=block))
    return pages or [PageText(page=1, text="\n".join(paragraphs))]


def _extract_pptx(path: Path) -> list[PageText]:
    from pptx import Presentation

    prs = Presentation(str(path))
    pages: list[PageText] = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
        text = "\n".join(parts).strip()
        if text:
            pages.append(PageText(page=idx, text=text))
    return pages


def _extract_txt(path: Path) -> list[PageText]:
    text = path.read_text(encoding="utf-8", errors="ignore")
    if not text.strip():
        return []
    # Split into pseudo-pages of ~3000 chars.
    page_size = 3000
    pages: list[PageText] = []
    for i in range(0, len(text), page_size):
        pages.append(PageText(page=(i // page_size) + 1, text=text[i : i + page_size]))
    return pages


EXTRACTORS = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
    ".txt": _extract_txt,
    ".md": _extract_txt,
    ".png": _extract_image,
    ".jpg": _extract_image,
    ".jpeg": _extract_image,
    ".webp": _extract_image,
    ".bmp": _extract_image,
    ".tiff": _extract_image,
}


def extract_pages(path: Path) -> list[PageText]:
    ext = path.suffix.lower()
    if ext not in EXTRACTORS:
        raise ValueError(f"Unsupported file type: {ext}")
    return EXTRACTORS[ext](path)


def chunk_pages(pages: list[PageText], chunk_size: int, chunk_overlap: int) -> list[Chunk]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", ". ", " ", ""],
    )
    chunks: list[Chunk] = []
    for pg in pages:
        for piece in splitter.split_text(pg.text):
            piece = piece.strip()
            if piece:
                chunks.append(Chunk(page=pg.page, text=piece))
    return chunks


def supported_extensions() -> set[str]:
    return set(EXTRACTORS.keys())
